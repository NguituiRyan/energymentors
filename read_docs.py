import sys
import os

try:
    from docx import Document
    from pptx import Presentation
except ImportError:
    pass

def extract():
    with open("c:/Users/phant/Desktop/Angawatch/content_clean.txt", "w", encoding="utf-8") as f:
        f.write("--- Extracted from DOCX ---\n")
        try:
            doc = Document("c:/Users/phant/Desktop/Angawatch/Angawatch_Complete_FIXED.docx")
            for para in doc.paragraphs:
                if para.text.strip():
                    f.write(para.text + "\n")
        except Exception as e:
            f.write(f"Error {e}\n")
        
        f.write("\n--- Extracted from PPTX ---\n")
        try:
            prs = Presentation("c:/Users/phant/Desktop/Angawatch/Angawatch_Template_Deck.pptx.pptx")
            for i, slide in enumerate(prs.slides):
                f.write(f"\nSlide {i+1}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        f.write(shape.text + "\n")
        except Exception as e:
            f.write(f"Error {e}\n")

if __name__ == "__main__":
    extract()
