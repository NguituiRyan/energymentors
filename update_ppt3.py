import os
import requests
import io
from pptx import Presentation
from pptx.util import Inches

# Use the presentation that has already had its fonts correctly shrunk
pptx_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Final_Deck.pptx"
output_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Final_Deck.pptx"
prs = Presentation(pptx_path)

images = {
    0: "1523805009346-ef9aca010fac",
    2: "1531123897727-8f129e1bfa82",
    3: "1542601906990-b4d3fb778b09",
    4: "1579621970163-a4f1092e0df4",
    10: "1541888081614-b1ecccf4183f",
    14: "1503387762-592deb58ef4e",
    15: "1503387762-592deb58ef4e",
    18: "1509391366360-2e959784a276",
    19: "1592982537443-d1e432f83d9e",
    24: "1585868202570-5cbab3357fb7",
    33: "1558981414-d03541c8d0cd",
    39: "1518770660439-4636190af475",
    43: "1550751827-4bd374c3f58b",
    46: "1551288049-bebda4e38f71",
    64: "1449844908441-8829872d2607"
}

def add_image_to_slide(slide, photo_id):
    try:
        url = f"https://images.unsplash.com/photo-{photo_id}?w=800&auto=format&fit=crop"
        response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=10)
        
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide.shapes.add_picture(image_stream, left, top, height=height)
            print(f"Added Unsplash photo {photo_id} to a slide.")
        else:
            print(f"Failed to download image {photo_id}, status code: {response.status_code}")
    except Exception as e:
        print(f"Failed to add image {photo_id}: {e}")

for slide_index, img_src in images.items():
    if slide_index < len(prs.slides):
        slide = prs.slides[slide_index]
        add_image_to_slide(slide, img_src)

prs.save(output_path)
print(f"Fully updated presentation with Unsplash photos saved to {output_path}")
