import os
import requests
import io
from pptx import Presentation
from pptx.util import Inches

# Load presentation
pptx_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Template_Deck.pptx.pptx"
output_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Optimized_Deck.pptx"
prs = Presentation(pptx_path)

# Image mapping (Slide index is 0-based, so Slide 1 is index 0)
images = {
    0: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\community_overview_1776103674674.png",
    2: "https://images.unsplash.com/photo-1531123897727-8f129e1bfa82", # Slide 3 (Stock)
    10: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\community_overview_1776103674674.png", # Slide 11
    14: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\ekeja_courtyard_1776103692120.png", # Slide 15
    15: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\cseb_construction_1776103707397.png", # Slide 16
    18: "https://images.unsplash.com/photo-1509391366360-2e959784a276", # Slide 19 (Stock)
    19: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\yattafarm_canal_1776103723892.png", # Slide 20
    24: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\mpesa_smart_meter_1776103738820.png", # Slide 25
    33: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\battery_swap_station_1776103754225.png", # Slide 34
    39: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\sensor_mesh_network_1776103772400.png", # Slide 40
    43: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\edge_inference_siren_1776103788666.png", # Slide 44
    46: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\ecohub_display_1776103806929.png", # Slide 47
    64: "https://images.unsplash.com/photo-1449844908441-8829872d2607", # Slide 65 (Stock)
}

def add_image_to_slide(slide, image_source):
    try:
        if image_source.startswith("http"):
            response = requests.get(image_source)
            image_stream = io.BytesIO(response.content)
        else:
            image_stream = image_source
        
        # Add to the right side of the slide
        # Assuming standard 16:9 ratio (10 inches wide, 5.625 inches high)
        left = Inches(5.5)
        top = Inches(1.5)
        height = Inches(3.5)
        slide.shapes.add_picture(image_stream, left, top, height=height)
        print(f"Added image to a slide.")
    except Exception as e:
        print(f"Failed to add image: {e}")

for slide_index, img_src in images.items():
    if slide_index < len(prs.slides):
        slide = prs.slides[slide_index]
        add_image_to_slide(slide, img_src)

prs.save(output_path)
print(f"Optimized presentation saved to {output_path}")
