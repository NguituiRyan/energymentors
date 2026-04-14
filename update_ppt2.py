import os
import requests
import io
from pptx import Presentation
from pptx.util import Inches

pptx_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Template_Deck.pptx.pptx"
output_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Optimized_Deck_Realistic.pptx"
prs = Presentation(pptx_path)

images = {
    0: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_community_1776105415490.png",
    2: "https://images.unsplash.com/photo-1531123897727-8f129e1bfa82", 
    10: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_community_1776105415490.png", 
    14: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_courtyard_1776105435350.png", 
    15: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_cseb_construction_1776105451321.png", 
    18: "https://images.unsplash.com/photo-1509391366360-2e959784a276", 
    19: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_yattafarm_1776105467226.png", 
    24: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_mpesa_meter_1776105484217.png", 
    33: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_boda_1776105501518.png", 
    39: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_sensor_mesh_1776105519165.png", 
    43: r"C:\Users\phant\.gemini\antigravity\brain\1a6d7c57-eb2d-4c47-81ea-b22a12306584\real_edge_siren_1776105538371.png", 
    46: "https://images.unsplash.com/photo-1551288049-bebda4e38f71", # Replaced last AI generation due to API quota
    64: "https://images.unsplash.com/photo-1449844908441-8829872d2607", 
}

def add_image_to_slide(slide, image_source):
    try:
        if image_source.startswith("http"):
            response = requests.get(image_source)
            image_stream = io.BytesIO(response.content)
        else:
            image_stream = image_source
        
        left = Inches(5.5)
        top = Inches(1.5)
        height = Inches(3.5)
        slide.shapes.add_picture(image_stream, left, top, height=height)
        print(f"Added image to a slide.")
    except Exception as e:
        print(f"Failed to add image: {e}")

for slide_index, img_src in images.items():
    if slide_index < len(prs.slides):
        slide = prs.slides[slide_index]
        add_image_to_slide(slide, img_src)

prs.save(output_path)
print(f"Optimized presentation saved to {output_path}")
