import io
import requests
from pptx import Presentation
from pptx.util import Inches

pptx_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Final_Deck.pptx"
output_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Final_Deck_Stock.pptx"
prs = Presentation(pptx_path)

# Working Unsplash photos
unsplash_images = {
    3: "1542601906990-b4d3fb778b09",
    14: "1503387762-592deb58ef4e",
    15: "1503387762-592deb58ef4e",
    18: "1509391366360-2e959784a276",
    39: "1518770660439-4636190af475",
    43: "1550751827-4bd374c3f58b",
    46: "1551288049-bebda4e38f71",
    64: "1449844908441-8829872d2607"
}

# Thematic dynamic stock photos for the ones that 404'd
lorem_images = {
    0: "africa,landscape",
    2: "african,woman",
    4: "african,business",
    10: "brick,building",
    19: "african,farming",
    24: "smartphone,hands",
    33: "scooter,charging"
}

def add_image(slide, url):
    try:
        response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'}, timeout=15)
        if response.status_code == 200:
            image_stream = io.BytesIO(response.content)
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide.shapes.add_picture(image_stream, left, top, height=height)
            print(f"Success: {url}")
        else:
            print(f"Failed status {response.status_code}: {url}")
    except Exception as e:
        print(f"Error {url}: {e}")

for slide_index, uid in unsplash_images.items():
    if slide_index < len(prs.slides):
        add_image(prs.slides[slide_index], f"https://images.unsplash.com/photo-{uid}?w=800&auto=format&fit=crop")

for slide_index, keyword in lorem_images.items():
    if slide_index < len(prs.slides):
        add_image(prs.slides[slide_index], f"https://loremflickr.com/800/600/{keyword}")

prs.save(output_path)
print(f"Fully updated presentation saved to {output_path}")
