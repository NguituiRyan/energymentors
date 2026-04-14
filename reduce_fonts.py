import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt

pptx_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Optimized_Deck_Realistic.pptx"
output_path = r"c:\Users\phant\Desktop\Angawatch\Angawatch_Final_Deck.pptx"
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    # The user mentioned editing the first few slides themselves. 
    # Let's adjust font sizes mostly on sliding indices after 4, but we'll apply a gentle downscale to all.
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        size_pt = run.font.size.pt
                        # If extremely large (like 100+), shrink to 54
                        if size_pt > 65:
                            run.font.size = Pt(54.0)
                        # If large heading (50-65), shrink to 36
                        elif size_pt >= 45:
                            run.font.size = Pt(36.0)
                        # Medium text (32-45), shrink to 24
                        elif size_pt >= 32:
                            run.font.size = Pt(24.0)
                        # Keep smaller standard text (20-30) as is, unless >28
                        elif size_pt > 28:
                            run.font.size = Pt(22.0)

prs.save(output_path)
print(f"Saved optimized font presentation to {output_path}")
